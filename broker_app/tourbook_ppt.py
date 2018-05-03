from pptx import Presentation
from .models import TourBook


class TourBookPPT:
    def __init__(self, buffer, tour_book: TourBook):
        self.buffer = buffer
        self.tour_book = tour_book

    def generate_ppt(self):
        prs = Presentation()
        bullet_slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = self.tour_book.tour_title

        tf = body_shape.text_frame
        tf.text = 'Surveys for tour:'

        for survey in self.tour_book.surveys.all():
            p = tf.add_paragraph()
            p.text = survey.building_name
            p.level = 1

        prs.save(self.buffer)
        buffer_value = self.buffer.getvalue()
        self.buffer.close()
        return buffer_value
